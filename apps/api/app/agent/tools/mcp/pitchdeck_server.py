"""Local, no-dependency Pitch Deck MCP server.

Generates a styled .pptx presentation entirely in-process using
python-pptx (already a project dependency) and returns it as a
base64-encoded ``__OPENHUMAN_FILE__`` marker, which the formatter node
picks up and auto-attaches to the chat — the same mechanism the
built-in ``create_document`` tool uses.

No API key, no OAuth, no network calls, no cost. Runs as a local
subprocess over stdio, same pattern as ``visualization_server.py``.
"""

import base64
import json
from typing import List

from mcp.server.fastmcp import FastMCP

mcp = FastMCP("pitchdeck")

_FILE_MARKER_PREFIX = "__OPENHUMAN_FILE__"

_BG_COLOR = (0x1A, 0x1A, 0x2E)
_TEXT_COLOR = (0xF0, 0xF0, 0xFF)
_ACCENT_COLOR = (0x6C, 0x63, 0xFF)
_MUTED_COLOR = (0xAF, 0xAF, 0xD9)


def _file_marker(filename: str, content_type: str, data_b64: str, title: str = "") -> str:
    payload = json.dumps({
        "filename": filename,
        "content_type": content_type,
        "data": data_b64,
        "title": title or filename,
    })
    return f"{_FILE_MARKER_PREFIX}{payload}"


def _build_pptx(company_name: str, tagline: str, slides: List[dict]) -> bytes:
    from io import BytesIO
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    bg_color = RGBColor(*_BG_COLOR)
    text_color = RGBColor(*_TEXT_COLOR)
    accent = RGBColor(*_ACCENT_COLOR)
    muted = RGBColor(*_MUTED_COLOR)

    def _new_slide():
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        return slide

    # -- Title slide -----------------------------------------------------
    title_slide = _new_slide()
    box = title_slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = company_name or "Untitled Pitch"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER

    if tagline:
        sub_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = tagline
        sp.font.size = Pt(22)
        sp.font.color.rgb = muted
        sp.alignment = PP_ALIGN.CENTER

    # -- Content slides ----------------------------------------------------
    for raw in slides or []:
        heading = str(raw.get("heading") or raw.get("title") or "").strip()
        bullets = raw.get("bullets") or []

        slide = _new_slide()
        if heading:
            head_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
            htf = head_box.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.text = heading
            hp.font.size = Pt(32)
            hp.font.bold = True
            hp.font.color.rgb = accent

        y_start = Inches(1.7) if heading else Inches(0.6)
        body_box = slide.shapes.add_textbox(Inches(0.9), y_start, Inches(11.5), Inches(5.4))
        btf = body_box.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(bullets):
            bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            bp.text = f"•  {bullet}"
            bp.font.size = Pt(20)
            bp.font.color.rgb = text_color
            bp.space_after = Pt(14)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


@mcp.tool()
async def create_pitch_deck(
    company_name: str,
    slides: List[dict],
    tagline: str = "",
    filename: str = "pitch_deck.pptx",
) -> str:
    """Generate a startup pitch deck as a downloadable .pptx file.

    Builds a styled, dark-themed presentation with a title slide
    (company name + tagline) followed by one slide per entry in
    ``slides``. The finished file is automatically attached to the
    chat — no further action needed.

    A typical pitch deck covers: Problem, Solution, Market Size,
    Product, Business Model, Traction, Competition, Team, and The Ask
    — but any set of slides works.

    Args:
        company_name: Name of the company/product for the title slide.
        slides: List of slides, each a dict with:
            - "heading": slide title (e.g. "Problem", "Market Size")
            - "bullets": list of short bullet-point strings
        tagline: One-line tagline shown under the company name.
        filename: Output filename, should end in ".pptx".
    """
    try:
        raw = _build_pptx(company_name, tagline, slides)
        data_b64 = base64.b64encode(raw).decode("ascii")
        return _file_marker(
            filename,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            data_b64,
            title=filename,
        )
    except Exception as e:
        return f"Error creating pitch deck: {e}"


if __name__ == "__main__":
    mcp.run(transport="stdio")
